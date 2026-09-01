from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any

from agentic_system.services.course_memory import CourseMemoryStore


class DocumentIngestionService:
    """Simple ingestion pipeline for course files downloaded from NTULearn or related course portals."""

    def __init__(self, db_path: str | None = None):
        self.store = CourseMemoryStore(db_path)

    def ingest_file(
        self,
        *,
        course_code: str,
        semester: str,
        file_path: str | Path,
        title: str | None = None,
        week_number: int | None = None,
        due_date: str | None = None,
        source_url: str | None = None,
    ) -> dict[str, Any]:
        path = Path(file_path)
        text = self._extract_text(path)
        summary = self._summarize_text(text)
        record = {
            "course_code": course_code,
            "semester": semester,
            "title": title or path.stem,
            "file_type": path.suffix.lower().lstrip("."),
            "week_number": week_number,
            "due_date": due_date,
            "source_url": source_url or "",
            "extracted_text": text,
            "summary": summary,
        }
        self.store.add_document(**record)
        return record

    def _extract_text(self, path: Path) -> str:
        if path.suffix.lower() == ".pdf":
            return self._extract_pdf_text(path)
        if path.suffix.lower() in {".txt", ".md", ".csv"}:
            return path.read_text(encoding="utf-8", errors="ignore")
        if path.suffix.lower() in {".ppt", ".pptx"}:
            return self._extract_ppt_text(path)
        return path.read_text(encoding="utf-8", errors="ignore") if path.exists() else ""

    def _extract_pdf_text(self, path: Path) -> str:
        try:
            import fitz

            with fitz.open(path) as doc:
                chunks = [page.get_text("text") for page in doc]
            return "\n".join(chunks)
        except Exception:
            return ""

    def _extract_ppt_text(self, path: Path) -> str:
        try:
            import pptx

            presentation = pptx.Presentation(path)
            chunks = []
            for slide in presentation.slides:
                texts = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        texts.append(shape.text)
                if texts:
                    chunks.append("\n".join(texts))
            return "\n".join(chunks)
        except Exception:
            return ""

    def _summarize_text(self, text: str) -> str:
        cleaned = " ".join(text.split())
        if not cleaned:
            return "No text extracted."
        if len(cleaned) <= 200:
            return cleaned
        return cleaned[:200] + "..."

    def fingerprint(self, file_path: str | Path) -> str:
        path = Path(file_path)
        digest = hashlib.sha256()
        digest.update(path.read_bytes())
        return digest.hexdigest()
